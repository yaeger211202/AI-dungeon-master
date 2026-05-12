"""
Generate the AI Dungeon Master demo deck (.pptx) — 10 minute solo presentation.

Run:
    cd backend
    python make_slides.py
    open ai_dungeon_master_demo.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---- Theme ----
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
DEEP_NAVY = RGBColor(0x16, 0x21, 0x3E)
GOLD = RGBColor(0xF0, 0xA5, 0x00)
CREAM = RGBColor(0xE8, 0xE2, 0xCC)
RED = RGBColor(0xE0, 0x70, 0x70)
GREEN = RGBColor(0x70, 0xE0, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0xA0)


def set_bg(slide, color):
    """Solid-color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=CREAM, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, bullets, size=20, color=CREAM, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_notes(slide, notes_text):
    slide.notes_slide.notes_text_frame.text = notes_text


def add_footer(slide, slide_num, total):
    add_text(slide, 0.4, 7.0, 9, 0.3,
             f"AI Dungeon Master  ·  CSC 603 Capstone  ·  Slide {slide_num}/{total}",
             size=10, color=GRAY, align=PP_ALIGN.LEFT)


# ---- Build deck ----

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
TOTAL = 9


# ===== Slide 1: Title =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 1.8, 9, 1.2, "⚔️ AI Dungeon Master ⚔️",
         size=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 3.1, 9, 0.8,
         "Structured Generative AI for Coherent Interactive Storytelling",
         size=22, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.6, 9, 0.5,
         "Krish Adya  ·  Harry Kakadiya  ·  Vansh Singh",
         size=18, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.3, 9, 0.4,
         "CSC 603 — Generative AI  ·  Spring 2026  ·  San Francisco State University",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(s,
    "Hi everyone, I'm Krish. Today my teammates Harry, Vansh and I are presenting "
    "AI Dungeon Master — a Generative AI system that uses a Large Language Model "
    "to run a fantasy role-playing game. But this isn't really a game project — "
    "it's an experiment in making LLMs stay consistent over long interactions. "
    "Let me show you why that matters."
)


# ===== Slide 2: The Problem =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.5, 9, 0.8, "LLMs Are Amazing Storytellers...",
         size=36, bold=True, color=GOLD)
add_text(s, 0.5, 1.3, 9, 0.6, "...until you talk to them for more than a few turns.",
         size=20, color=CREAM)

add_text(s, 0.5, 2.3, 9, 0.5, "Vanilla LLM fails in interactive storytelling:",
         size=18, bold=True, color=GOLD)
add_bullets(s, 0.7, 2.9, 9, 4, [
    "Hallucinated items — gives you a Magic Sword you never picked up",
    "Teleporting players — you blink and you're in the Dragon's Lair",
    "Quest amnesia — quest text silently changes mid-game",
    "Health drift — restored to 100 HP for no reason, or drops with no fight",
    "Forgets context after only 3-4 turns of conversation",
], size=20)

add_text(s, 0.5, 6.3, 9, 0.5,
         "These break the illusion of a persistent game world.",
         size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_footer(s, 2, TOTAL)
add_notes(s,
    "LLMs like GPT-4 or Llama-3 are incredible at creative writing. "
    "But the moment you ask them to maintain a persistent world over many turns — "
    "track your inventory, remember your quest, keep your health consistent — they fall apart. "
    "They invent items. They teleport you across the map. They quietly swap your quest. "
    "And after just 3-4 turns, they forget what happened earlier. "
    "This is the problem we set out to fix."
)


# ===== Slide 3: Our Approach =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.5, 9, 0.8, "What If We Don't Trust the LLM Alone?",
         size=32, bold=True, color=GOLD)

add_text(s, 0.5, 1.7, 9, 0.5, "Our core idea:",
         size=20, bold=True, color=GOLD)
add_bullets(s, 0.7, 2.3, 9, 4, [
    "The LLM still tells the story — it's a creative narrator",
    "But the server holds the truth — independent state manager",
    "Every state change is justified by the narrative or it gets rejected",
    "Long memory comes from compressed summaries, not raw history",
], size=20)

add_text(s, 0.5, 5.3, 9, 0.5,
         "Two systems, same LLM:",
         size=18, bold=True, color=GOLD)
add_bullets(s, 0.7, 5.8, 9, 1.5, [
    "Baseline — pure free-form LLM (the control)",
    "Structured — LLM + state mgr + validator + memory (our system)",
], size=18, color=CREAM)
add_footer(s, 3, TOTAL)
add_notes(s,
    "Our approach is simple: treat the LLM like a creative writer, not a database. "
    "It tells the story, it writes the dialogue, it generates the choices. "
    "But the SERVER decides what's actually true. "
    "If the LLM says you found a magic sword, the validator checks: did the narrative actually describe finding a sword? "
    "If not — rejected. The change never persists. "
    "And we built two versions side-by-side: a baseline with no safeguards, and our structured system. "
    "You'll see them compete live in a minute."
)


# ===== Slide 4: Architecture =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.4, 9, 0.7, "System Architecture",
         size=32, bold=True, color=GOLD)

# Box drawing
boxes = [
    (0.5, 1.5, 2.2, 0.9, "Frontend\n(HTML + JS)", CREAM),
    (3.1, 1.5, 2.2, 0.9, "FastAPI\nBackend", CREAM),
    (5.7, 1.5, 2.2, 0.9, "Groq LLM\n(Llama-3.3-70B)", CREAM),
    (0.5, 3.0, 2.2, 0.9, "State Manager\n(per-session JSON)", GOLD),
    (3.1, 3.0, 2.2, 0.9, "Validator\n(7 rules)", GOLD),
    (5.7, 3.0, 2.2, 0.9, "Memory Mgr\n(every 6 turns)", GOLD),
    (0.5, 4.5, 7.4, 0.9, "Logs (JSONL) — every validation decision saved", DEEP_NAVY),
]
for x, y, w, h, label, fill_color in boxes:
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = DEEP_NAVY if fill_color != DEEP_NAVY else NAVY
    shp.line.color.rgb = fill_color if fill_color != DEEP_NAVY else GOLD
    shp.line.width = Pt(2)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = fill_color if fill_color != DEEP_NAVY else CREAM

add_text(s, 0.5, 5.7, 9, 0.5,
         "Gold boxes = our research contribution. White boxes = stock infrastructure.",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, 0.5, 6.3, 9, 0.5,
         "Per turn: Player → Backend → LLM proposes → Validator checks → State Mgr persists.",
         size=15, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
add_footer(s, 4, TOTAL)
add_notes(s,
    "Here's the architecture. The white boxes — frontend, backend, and the LLM itself — are off-the-shelf. "
    "Our research contribution is the three gold boxes plus the logging layer. "
    "Every turn flows through this pipeline: the player sends an action, "
    "we build a prompt with state and memory, the LLM proposes new state in JSON, "
    "the validator checks it against 7 rules, and only valid changes get persisted. "
    "Every decision the validator makes is logged so we can measure how often it catches the LLM."
)


# ===== Slide 5: LIVE DEMO transition =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 2.8, 9, 1.2, "🎲  LIVE DEMO  🎲",
         size=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.3, 9, 0.6, "Let's play.",
         size=28, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.4, 9, 0.5,
         "Watch for: state tracking · choice buttons · validator warnings · vs baseline contradictions",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 5, TOTAL)
add_notes(s,
    "DEMO SCRIPT (roughly 3-4 minutes):\n"
    "1. Open the frontend in browser. Click '🔄 New Game' in Structured mode.\n"
    "2. Play 3-4 turns of a coherent adventure. Use the choice buttons sometimes.\n"
    "   Point out the stats panel updating in real time (health, inventory, location).\n"
    "3. After a turn — if you happen to get a violation banner, point it out:\n"
    "   'See — the validator caught the LLM trying to change something it shouldn't.'\n"
    "4. CLICK 'Baseline (Free-form)' button. Click 🔄 New Game.\n"
    "   Show that the stats panel is now dimmed — no structured tracking.\n"
    "5. Play 2-3 baseline turns. Then ask a state-dependent question:\n"
    "   'What's in my inventory?' — baseline often invents new items.\n"
    "   'How much health do I have?' — baseline often pulls a number out of thin air.\n"
    "6. CLICK 'Structured' button. Show the structured story is preserved.\n"
    "   Ask the SAME inventory question — it gives the real list, every time.\n"
    "7. Wrap up: 'Same LLM. Same input. Different system architecture.'"
)


# ===== Slide 6: Validator + Memory (combined) =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.4, 9, 0.7, "What You Just Saw",
         size=32, bold=True, color=GOLD)

# Left column: Validator
add_text(s, 0.4, 1.4, 4.6, 0.5, "🛡️  Validator — 7 rules",
         size=20, bold=True, color=GOLD)
add_bullets(s, 0.5, 2.0, 4.6, 4.5, [
    "HEALTH_BOUNDS  (clamp 0-100)",
    "HEALTH_JUMP  (big delta needs combat / healing)",
    "INVENTORY_ADD_UNJUSTIFIED",
    "INVENTORY_REMOVE_UNJUSTIFIED",
    "LOCATION_TELEPORT  (needs movement verb)",
    "QUEST_DRIFT  (silent quest swaps)",
    "OUTPUT_SCHEMA  (missing JSON fields)",
], size=15, color=CREAM)

# Right column: Memory
add_text(s, 5.2, 1.4, 4.5, 0.5, "🧠  Memory Summarization",
         size=20, bold=True, color=GOLD)
add_bullets(s, 5.3, 2.0, 4.5, 4.5, [
    "Every 6 turns the LLM summarizes the chunk in 2-3 sentences",
    "Summaries are added to every future prompt",
    "Raw last 6 lines also included",
    "Result: stays coherent across 15-20+ turns",
    "No token blowup",
], size=15, color=CREAM)
add_footer(s, 6, TOTAL)
add_notes(s,
    "The two innovations under the hood. "
    "The Validator runs 7 rules over every LLM response. "
    "Each rule has a specific failure mode it watches for. "
    "For example: HEALTH_JUMP triggers when health changes by more than 40 in one turn — "
    "but only fires if the narrative doesn't mention combat or healing. "
    "If the LLM says 'an orc slashes you', a health drop is fine. "
    "If it says 'you look at the trees' and drops your health to 10, that's caught and reverted.\n\n"
    "On the right — memory. Without summarization, the LLM forgets anything past turn 3 or 4. "
    "Every 6 turns we ask it to compress that chunk into 2-3 sentences, "
    "and those summaries become part of the prompt forever after. "
    "That's how we hit our 15-20 turn coherence target."
)


# ===== Slide 7: Results =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.4, 9, 0.7, "Does It Actually Work?",
         size=32, bold=True, color=GOLD)

# Big number callouts
def big_stat(x, y, value, label, color=GOLD):
    add_text(s, x, y, 3, 1.1, value, size=64, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, y + 1.1, 3, 0.6, label, size=14, color=CREAM, align=PP_ALIGN.CENTER)

big_stat(0.5, 1.6, "66", "TURNS PLAYED (STRUCTURED)")
big_stat(3.5, 1.6, "2", "VIOLATIONS CAUGHT")
big_stat(6.5, 1.6, "3.0%", "VIOLATION RATE", color=GREEN)

# Target callout
add_text(s, 0.5, 4.4, 9, 0.6,
         "Proposal target: < 5% inconsistency rate  →  3.0%  ✅  PASS",
         size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Bullet observations
add_bullets(s, 0.5, 5.3, 9, 2, [
    "Both violations were LOCATION_TELEPORT — caught and reverted automatically",
    "Baseline transcripts show contradictions you can spot by eye (demo'd above)",
    "Every validator decision is logged — full audit trail in JSONL files",
], size=16)
add_footer(s, 7, TOTAL)
add_notes(s,
    "These are the real numbers from playing through the system. "
    "Over 66 turns of gameplay, the validator caught only 2 inconsistencies — "
    "a violation rate of just 3 percent. "
    "Our proposal set a target of under 5 percent — we beat it. "
    "Both of the violations we caught were LOCATION_TELEPORT — the LLM tried to silently change the player's location. "
    "Both were reverted automatically without breaking the game. "
    "And critically, the baseline system, with no validation, can't even tell you what its own state is. "
    "I asked it what I had in my inventory — it just invented things."
)


# ===== Slide 8: Takeaways =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 0.4, 9, 0.7, "Takeaways",
         size=32, bold=True, color=GOLD)
add_bullets(s, 0.7, 1.5, 9, 4, [
    "LLMs are creative narrators, not databases — don't make them be both",
    "Structured output (JSON mode) + server-side validation is cheap and effective",
    "Memory summarization is simple but unlocks 15-20+ turn coherence",
    "Open-source 8B/70B models work fine for this — no need for GPT-4",
], size=20)

add_text(s, 0.5, 5.3, 9, 0.6, "Future work:",
         size=20, bold=True, color=GOLD)
add_bullets(s, 0.7, 5.9, 9, 1.5, [
    "LLM-as-judge for semantic validation (catch subtler hallucinations)",
    "RAG over a world bible for lore consistency",
    "Multiplayer / shared-world support",
], size=16)
add_footer(s, 8, TOTAL)
add_notes(s,
    "Three things we learned. First — and this is the big one — separate creative work from bookkeeping. "
    "The LLM is great at telling stories. It's bad at maintaining a database. "
    "Don't make it do both. Let it propose, and let your code decide what's true.\n\n"
    "Second — JSON mode plus a hundred lines of validation code goes a really long way. "
    "We didn't need fine-tuning or a giant model. Llama-3.3-70B over Groq's free tier was enough.\n\n"
    "Third — for long sessions, memory compression is a must. Without summarization the LLM forgets fast.\n\n"
    "If we kept going: an LLM-as-judge could catch subtler issues than our keyword rules can. "
    "And a retrieval system over a 'world bible' would help with lore consistency across very long campaigns."
)


# ===== Slide 9: Thank You =====
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_text(s, 0.5, 2.0, 9, 1.0, "Thank You",
         size=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 3.4, 9, 0.6, "Questions?",
         size=32, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.7, 9, 0.5,
         "Krish Adya  ·  Harry Kakadiya  ·  Vansh Singh",
         size=18, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.4, 9, 0.5,
         "CSC 603 — Generative AI Capstone  ·  Spring 2026",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 9, TOTAL)
add_notes(s,
    "Thank you! Happy to answer questions about the architecture, the validator rules, "
    "the memory system, or anything you saw in the demo."
)


# ---- Save ----
out = Path(__file__).parent / "ai_dungeon_master_demo.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
print(f"   {TOTAL} slides — speaker notes included on every slide.")
